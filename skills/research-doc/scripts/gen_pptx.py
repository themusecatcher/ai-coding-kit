#!/usr/bin/env python3
"""
PPT 调研报告生成器
用法：python3 gen_pptx.py --input content.json --output report.pptx
"""

import argparse
import json
import sys
import os
from datetime import datetime

def check_deps():
    try:
        import pptx
    except ImportError:
        print("❌ 缺少 python-pptx，请执行: pip3 install --user python-pptx")
        sys.exit(1)

def create_pptx(data: dict, output_path: str):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── 配色方案 ──
    COLOR_PRIMARY = RGBColor(0x1A, 0x73, 0xE8)    # 蓝色主色
    COLOR_DARK = RGBColor(0x20, 0x2A, 0x44)        # 深色文字
    COLOR_SECONDARY = RGBColor(0x5F, 0x6B, 0x7C)   # 次级文字
    COLOR_ACCENT = RGBColor(0x00, 0xC4, 0xB4)       # 强调色
    COLOR_BG_LIGHT = RGBColor(0xF7, 0xF8, 0xFA)    # 浅灰背景
    COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    def add_bg(slide, color):
        """为幻灯片添加纯色背景"""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_rect(slide, left, top, width, height, color, transparency=0):
        """添加矩形色块"""
        from pptx.util import Emu
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def set_text(tf, text, size=18, color=COLOR_DARK, bold=False, alignment=PP_ALIGN.LEFT):
        """设置文本框内容"""
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = alignment

    # ════════════════════════════════════════
    # 封面页
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    add_bg(slide, COLOR_PRIMARY)

    # 左侧装饰条
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), slide_height, COLOR_ACCENT)

    # 标题
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10), Inches(1.5))
    set_text(txBox.text_frame, data.get("title", "调研报告"), size=42, color=COLOR_WHITE, bold=True)

    # 副标题
    if data.get("subtitle"):
        txBox = slide.shapes.add_textbox(Inches(1.2), Inches(3.6), Inches(10), Inches(0.8))
        set_text(txBox.text_frame, data["subtitle"], size=22, color=RGBColor(0xBB, 0xDE, 0xFB))

    # 作者和日期
    meta_parts = []
    if data.get("author"):
        meta_parts.append(data["author"])
    meta_parts.append(data.get("date", datetime.now().strftime("%Y-%m-%d")))
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(5.5), Inches(10), Inches(0.5))
    set_text(txBox.text_frame, " | ".join(meta_parts), size=16, color=RGBColor(0xBB, 0xDE, 0xFB))

    # ════════════════════════════════════════
    # 目录页（当章节 >= 3 时生成）
    # ════════════════════════════════════════
    sections = data.get("sections", [])
    if len(sections) >= 3:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide, COLOR_WHITE)

        # 页面标题
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(5), Inches(0.8))
        set_text(txBox.text_frame, "目录", size=32, color=COLOR_PRIMARY, bold=True)

        # 分隔线
        add_rect(slide, Inches(1.0), Inches(1.4), Inches(2), Inches(0.06), COLOR_ACCENT)

        # 目录条目
        for i, sec in enumerate(sections):
            y = Inches(2.0 + i * 0.65)
            if y > Inches(6.5):
                break

            # 序号圆点
            num_box = slide.shapes.add_textbox(Inches(1.0), y, Inches(0.5), Inches(0.5))
            set_text(num_box.text_frame, f"{i + 1:02d}", size=16, color=COLOR_PRIMARY, bold=True)

            # 标题
            title_box = slide.shapes.add_textbox(Inches(1.8), y, Inches(9), Inches(0.5))
            set_text(title_box.text_frame, sec.get("title", f"章节 {i + 1}"), size=18, color=COLOR_DARK)

    # ════════════════════════════════════════
    # 内容页
    # ════════════════════════════════════════
    for i, section in enumerate(sections):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide, COLOR_WHITE)

        sec_type = section.get("type", "text")

        # 顶部色带
        add_rect(slide, Inches(0), Inches(0), slide_width, Inches(0.06), COLOR_PRIMARY)

        # 章节编号 + 标题
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
        tf = txBox.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        from pptx.util import Pt
        run_num = p.add_run()
        run_num.text = f"{i + 1:02d}  "
        run_num.font.size = Pt(14)
        run_num.font.color.rgb = COLOR_ACCENT
        run_num.font.bold = True

        run_title = p.add_run()
        run_title.text = section.get("title", f"章节 {i + 1}")
        run_title.font.size = Pt(28)
        run_title.font.color.rgb = COLOR_DARK
        run_title.font.bold = True

        # 标题下分隔线
        add_rect(slide, Inches(0.8), Inches(1.25), Inches(1.5), Inches(0.04), COLOR_ACCENT)

        content = section.get("content", "")

        if sec_type == "table" and section.get("data"):
            # ── 表格类型 ──
            table_data = section["data"]
            headers = table_data.get("headers", [])
            rows = table_data.get("rows", [])
            if headers and rows:
                cols = len(headers)
                row_count = min(len(rows) + 1, 12)  # 限制行数
                table_width = Inches(11.5)
                col_width = table_width // cols

                table_shape = slide.shapes.add_table(
                    row_count, cols,
                    Inches(0.8), Inches(1.6),
                    table_width, Inches(0.45 * row_count)
                )
                table = table_shape.table

                # 表头
                for j, h in enumerate(headers):
                    cell = table.cell(0, j)
                    cell.text = str(h)
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(14)
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = COLOR_WHITE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLOR_PRIMARY

                # 数据行
                for r_idx, row in enumerate(rows[:row_count - 1]):
                    for c_idx, val in enumerate(row):
                        if c_idx < cols:
                            cell = table.cell(r_idx + 1, c_idx)
                            cell.text = str(val)
                            for paragraph in cell.text_frame.paragraphs:
                                paragraph.font.size = Pt(13)
                                paragraph.font.color.rgb = COLOR_DARK
                            if r_idx % 2 == 1:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = COLOR_BG_LIGHT

        elif sec_type == "bullet" or "\n- " in content or "\n• " in content:
            # ── 要点列表 ──
            lines = content.replace("• ", "- ").split("\n")
            bullets = [l.lstrip("- ").strip() for l in lines if l.strip().startswith("-") or l.strip()]
            if not any(l.strip().startswith("-") for l in lines):
                bullets = [l.strip() for l in lines if l.strip()]

            y_start = Inches(1.6)
            for b_idx, bullet in enumerate(bullets[:10]):
                if not bullet:
                    continue
                y = y_start + Inches(b_idx * 0.55)
                if y > Inches(6.5):
                    break

                # 圆点
                dot = slide.shapes.add_textbox(Inches(0.8), y, Inches(0.3), Inches(0.4))
                set_text(dot.text_frame, "●", size=10, color=COLOR_ACCENT)

                # 文本
                txt = slide.shapes.add_textbox(Inches(1.3), y, Inches(10.5), Inches(0.5))
                set_text(txt.text_frame, bullet, size=16, color=COLOR_DARK)

        else:
            # ── 文本类型 ──
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.2))
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.clear()

            paragraphs = content.split("\n")
            for p_idx, para_text in enumerate(paragraphs):
                if not para_text.strip():
                    continue
                if p_idx == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = para_text.strip()
                p.font.size = Pt(16)
                p.font.color.rgb = COLOR_DARK
                p.space_after = Pt(8)
                p.line_spacing = Pt(24)

        # 页码
        page_box = slide.shapes.add_textbox(
            Inches(12.0), Inches(6.9), Inches(1), Inches(0.4)
        )
        # 封面 + 目录不计入页码
        page_num = i + 1 + (2 if len(sections) >= 3 else 1)
        set_text(page_box.text_frame, str(page_num), size=11, color=COLOR_SECONDARY, alignment=PP_ALIGN.RIGHT)

    # ════════════════════════════════════════
    # 总结页
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, COLOR_PRIMARY)
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), slide_height, COLOR_ACCENT)

    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(2.8), Inches(10), Inches(1.2))
    set_text(txBox.text_frame, "Thank You", size=44, color=COLOR_WHITE, bold=True, alignment=PP_ALIGN.LEFT)

    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(4.2), Inches(10), Inches(0.6))
    set_text(txBox.text_frame, data.get("date", datetime.now().strftime("%Y-%m-%d")),
             size=18, color=RGBColor(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.LEFT)

    # 保存
    prs.save(output_path)
    print(f"✅ PPT 已生成：{output_path}")
    print(f"   共 {len(prs.slides)} 页")


def main():
    parser = argparse.ArgumentParser(description="PPT 调研报告生成器")
    parser.add_argument("--input", required=True, help="内容 JSON 文件路径")
    parser.add_argument("--output", required=True, help="输出 PPTX 文件路径")
    args = parser.parse_args()

    check_deps()

    with open(args.input, "r", encoding="utf-8") as f:
        data = json.load(f)

    # 确保输出目录存在
    output_dir = os.path.dirname(args.output)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)

    create_pptx(data, args.output)


if __name__ == "__main__":
    main()
